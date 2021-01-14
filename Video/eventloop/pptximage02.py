import os 
import grp
import pwd
from pptx import Presentation
from pptx.util import Inches
import datetime
from datetime import date

def next_weekday(d, weekday):
    days_ahead = weekday - d.weekday()
    if days_ahead <= 0: # Target day already happened this week
        days_ahead += 7
    return (d + datetime.timedelta(days_ahead)).strftime("%Y%m%d")

#today = date.today()
weekno = datetime.datetime.today().weekday()
d = date.today()
if weekno == 6:
	next_sunday = (d).strftime("%Y%m%d")
if weekno<6:
	next_sunday = next_weekday(d, 6) # 0 = Monday, 1=Tuesday, 2=Wednesday...
 
prs = Presentation()
factor = .5
prs.slide_height = Inches(11.25 * factor) 
prs.slide_width = Inches(20 * factor)

if not os.path.exists('sonntag'):
    os.makedirs('sonntag')
def change_permissions_recursive(path, mode):
    for root, dirs, files in os.walk(path, topdown=False):
        for dir in [os.path.join(root,d) for d in dirs]:
            os.chmod(dir, mode)
    for file in [os.path.join(root, f) for f in files]:
            os.chmod(file, mode)
#change_permissions_recursive('sonntag', 0o777)
 
for root, dirs , filenames in os.walk("./{0}/".format(next_sunday)):
  valid_images = [".jpg",".png"]
  for file in sorted(filenames):
    ext = os.path.splitext(file)[1]
    if ext.lower() in valid_images:
      image_path=os.path.join(root,file)
      blank_slide_layout = prs.slide_layouts[6]
      slide = prs.slides.add_slide(blank_slide_layout)
      top = Inches(0)
      left = Inches(0)
      height = Inches(5.67)
      pic = slide.shapes.add_picture(image_path,left,top, height=height)
prs.save("./sonntag/{0}.pptx".format(next_sunday))
os.chmod("./sonntag/{0}.pptx".format(next_sunday), 0o0777)
uid = pwd.getpwnam("keith").pw_uid
gid = grp.getgrnam("users").gr_gid
os.chown("./sonntag/{0}.pptx".format(next_sunday), uid, gid)
